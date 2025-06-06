"""
Content providers for extracting text from various file formats.
"""
import tempfile
from pathlib import Path
from abc import ABC, abstractmethod
from urllib.parse import urlparse
import requests
import fitz  # PyMuPDF
from percolate import logger

# Try importing docx libraries
try:
    from docx import Document
    has_docx = True
except ImportError:
    has_docx = False
    logger.warning("python-docx not installed, DOCX support limited")

try:
    import mammoth
    has_mammoth = True
except ImportError:
    has_mammoth = False
    logger.warning("mammoth not installed, DOCX support limited")

try:
    import html2text
    has_html2text = True
except ImportError:
    has_html2text = False
    logger.warning("html2text not installed, DOCX HTML conversion not available")


def is_url(uri: str) -> bool:
    parsed = urlparse(uri)
    return parsed.scheme in ("http", "https")


def resolve_path_or_download(uri: str) -> Path:
    if Path(uri).exists():
        return Path(uri)

    if is_url(uri):
        response = requests.get(uri)
        response.raise_for_status()
        suffix = Path(urlparse(uri).path).suffix
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp.write(response.content)
        tmp.close()
        return Path(tmp.name)

    raise FileNotFoundError(f"Cannot resolve URI: {uri}")


class BaseContentProvider(ABC):
    @abstractmethod
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        """Extract text from a file.
        
        Args:
            uri: File path or URL
            enriched: If True, use advanced processing (LLM analysis, etc.)
                     If False, use basic text extraction
        """
        ...


class PDFContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        if not enriched:
            # Raw mode: simple text extraction
            with fitz.open(str(path)) as doc:
                return "\n".join(page.get_text() for page in doc)
        else:
            # Enriched mode: TODO - convert pages to images and use LLM
            logger.warning("PDF enriched mode not yet implemented, falling back to raw")
            with fitz.open(str(path)) as doc:
                return "\n".join(page.get_text() for page in doc)


class DefaultContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        text = path.read_text()
        
        if enriched:
            # Enriched mode: TODO - use LLM to interpret/summarize content
            logger.warning("Default enriched mode not yet implemented, falling back to raw")
        
        return text


class DOCXContentProvider(BaseContentProvider):
    """Content provider for Microsoft Word documents."""
    
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        """Extract text from a DOCX file."""
        path = resolve_path_or_download(uri)
        
        # If no libraries available, fall back to basic text extraction
        if not has_docx and not has_mammoth:
            logger.warning("No DOCX libraries available, falling back to simple text extraction")
            return path.read_text(errors='ignore')
        
        try:
            # First try with python-docx for simple text extraction
            if has_docx:
                doc = Document(str(path))
                paragraphs = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        paragraphs.append(paragraph.text)
                
                # If we got text, return it
                if paragraphs:
                    return '\n\n'.join(paragraphs)
            
            # Fallback to mammoth for more complex documents
            if has_mammoth:
                logger.info("Using mammoth for DOCX extraction")
                with open(str(path), "rb") as docx_file:
                    result = mammoth.convert_to_markdown(docx_file)
                    
                    if result.messages:
                        for message in result.messages:
                            logger.warning(f"DOCX conversion warning: {message}")
                    
                    return result.value
                    
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {e}")
            # Final fallback - try mammoth HTML conversion if available
            if has_mammoth and has_html2text:
                try:
                    with open(str(path), "rb") as docx_file:
                        result = mammoth.convert_to_html(docx_file)
                        h = html2text.HTML2Text()
                        h.ignore_links = False
                        return h.handle(result.value)
                except Exception as e2:
                    logger.error(f"Failed all DOCX extraction methods: {e2}")
            
            # Ultimate fallback
            return path.read_text(errors='ignore')


# New provider classes for additional formats
class HTMLContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        if not enriched:
            # Raw mode: strip HTML tags
            try:
                if has_html2text:
                    h = html2text.HTML2Text()
                    h.ignore_links = True
                    return h.handle(path.read_text())
                else:
                    import re
                    html_content = path.read_text()
                    return re.sub(r'<[^>]+>', '', html_content)
            except Exception as e:
                logger.error(f"HTML parsing failed: {e}")
                return path.read_text()
        else:
            # Enriched mode: TODO - use LLM to analyze structure and content
            logger.warning("HTML enriched mode not yet implemented, falling back to raw")
            return self.extract_text(uri, enriched=False)


class MarkdownContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        text = path.read_text()
        
        if enriched:
            # Enriched mode: TODO - parse markdown structure, extract metadata
            logger.warning("Markdown enriched mode not yet implemented, falling back to raw")
        
        return text


class XLSXContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        try:
            import pandas as pd
            df = pd.read_excel(str(path), sheet_name=None)  # Read all sheets
            
            if not enriched:
                # Raw mode: concatenate all cell values
                all_text = []
                for sheet_name, sheet_df in df.items():
                    all_text.append(f"Sheet: {sheet_name}\n")
                    all_text.append(sheet_df.to_string(index=False))
                return "\n\n".join(all_text)
            else:
                # Enriched mode: TODO - analyze data structure, detect patterns
                logger.warning("XLSX enriched mode not yet implemented, falling back to raw")
                return self.extract_text(uri, enriched=False)
                
        except ImportError:
            logger.warning("pandas not available, falling back to default provider")
            return DefaultContentProvider().extract_text(uri, enriched=enriched)
        except Exception as e:
            logger.error(f"XLSX parsing failed: {e}")
            return f"Error reading XLSX file: {e}"


class WAVContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        if not enriched:
            # Raw mode: basic file info
            try:
                import wave
                with wave.open(str(path), 'rb') as wav_file:
                    frames = wav_file.getnframes()
                    sample_rate = wav_file.getframerate()
                    duration = frames / sample_rate
                    return f"Audio file: {path.name}\nDuration: {duration:.2f} seconds\nSample rate: {sample_rate} Hz\nFrames: {frames}"
            except Exception as e:
                return f"Audio file: {path.name}\nError reading file: {e}"
        else:
            # Enriched mode: TODO - speech-to-text transcription
            logger.warning("WAV enriched mode (speech-to-text) not yet implemented, falling back to raw")
            return self.extract_text(uri, enriched=False)


class PPTXContentProvider(BaseContentProvider):
    def extract_text(self, uri: str, enriched: bool = False) -> str:
        path = resolve_path_or_download(uri)
        
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            
            if not enriched:
                # Raw mode: extract all text from slides
                text_runs = []
                for slide_num, slide in enumerate(prs.slides, 1):
                    text_runs.append(f"Slide {slide_num}:")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                return "\n\n".join(text_runs)
            else:
                # Enriched mode: TODO - analyze slide structure, extract images
                logger.warning("PPTX enriched mode not yet implemented, falling back to raw")
                return self.extract_text(uri, enriched=False)
                
        except ImportError:
            logger.warning("python-pptx not available, falling back to default provider")
            return DefaultContentProvider().extract_text(uri, enriched=enriched)
        except Exception as e:
            logger.error(f"PPTX parsing failed: {e}")
            return f"Error reading PPTX file: {e}"


content_providers = {
    ".pdf": PDFContentProvider(),
    ".docx": DOCXContentProvider(),
    ".doc": DOCXContentProvider(),  # Will handle old doc format too
    ".txt": DefaultContentProvider(),
    ".html": HTMLContentProvider(),
    ".htm": HTMLContentProvider(),
    ".md": MarkdownContentProvider(),
    ".markdown": MarkdownContentProvider(),
    ".xlsx": XLSXContentProvider(),
    ".xls": XLSXContentProvider(),
    ".wav": WAVContentProvider(),
    ".pptx": PPTXContentProvider(),
    ".ppt": PPTXContentProvider(),
}

default_provider = DefaultContentProvider()


def get_content_provider_for_uri(uri: str) -> BaseContentProvider:
    """Get the appropriate content provider for a given URI."""
    suffix = Path(urlparse(uri).path).suffix.lower()
    return content_providers.get(suffix, default_provider)